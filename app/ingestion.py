import json
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Literal, Optional, Iterable
import hashlib
import chardet

# Document loaders
from langchain_community.document_loaders import (
    PyPDFLoader, TextLoader, CSVLoader
)
from langchain_core.documents import Document
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    MarkdownTextSplitter,
    PythonCodeTextSplitter,
    Language
)
from langchain_openai import OpenAIEmbeddings
from langchain_pinecone import PineconeVectorStore

# File processing libraries
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup
import markdown
from striprtf.striprtf import rtf_to_text

from pinecone import Pinecone
from pinecone.exceptions import PineconeApiException

from app.settings import settings

Domain = Literal["therapy", "health_fitness", "literature"]

# File type mappings for code files
CODE_EXTENSIONS = {
    '.py': Language.PYTHON,
    '.js': Language.JS,
    '.jsx': Language.JS,
    '.ts': Language.TS,
    '.tsx': Language.TS,
    '.java': Language.JAVA,
    '.cpp': Language.CPP,
    '.c': Language.C,
    '.cs': Language.CSHARP,
    '.go': Language.GO,
    '.php': Language.PHP,
    '.rb': Language.RUBY,
    '.rs': Language.RUST,
    '.scala': Language.SCALA,
    '.swift': Language.SWIFT,
    '.kt': Language.KOTLIN,
    '.html': Language.HTML,
    # Unsupported by langchain Language enum, will use generic text processing
    '.r': None,
    '.sql': None,
    '.css': None,
    '.scss': None,
    '.less': None,
}

def _index_name(domain: Domain) -> str:
    if domain == "therapy":
        return settings.INDEX_THERAPY
    if domain == "health_fitness":
        return settings.INDEX_HEALTH
    return settings.INDEX_LITERATURE


# ---------- Stable IDs to avoid duplicates ----------

def _stable_doc_id(d: Document) -> str:
    """Deterministic ID from content + key metadata so re-ingesting
    the same file doesn't duplicate vectors."""
    src = str(d.metadata.get("source", ""))
    page = str(d.metadata.get("page", ""))
    sheet = str(d.metadata.get("sheet", ""))
    row = str(d.metadata.get("row_id", ""))
    slide = str(d.metadata.get("slide", ""))
    text = d.page_content.strip()
    key = f"{src}|{page}|{sheet}|{row}|{slide}|{hashlib.sha256(text.encode('utf-8')).hexdigest()}"
    return hashlib.sha256(key.encode("utf-8")).hexdigest()[:32]


def _batched(seq: List[str], size: int) -> Iterable[List[str]]:
    for i in range(0, len(seq), size):
        yield seq[i : i + size]


def _fetch_existing_ids(
    index, ids: List[str], namespace: str | None, initial_batch: int = 100
) -> set[str]:
    """Fetch existing vectors' IDs in batches to avoid 414 Request-URI Too Large."""
    existing: set[str] = set()
    batch_size = initial_batch
    i = 0
    while i < len(ids):
        batch = ids[i : i + batch_size]
        try:
            resp = index.fetch(ids=batch, namespace=(namespace or ""))
            existing.update((resp.vectors or {}).keys())
            i += batch_size
        except PineconeApiException as e:
            if getattr(e, "status", None) == 414 or "414" in str(e):
                if batch_size <= 10:
                    raise
                batch_size //= 2
            else:
                raise
    return existing


def _filter_new_docs(
    index_name: str,
    namespace: str | None,
    docs: List[Document]
) -> tuple[List[Document], List[str]]:
    """Ask Pinecone which IDs already exist; return only the new docs + their IDs."""
    ids = [_stable_doc_id(d) for d in docs]

    pc = Pinecone(api_key=settings.PINECONE_API_KEY)
    index = pc.Index(index_name)

    existing_ids = _fetch_existing_ids(index, ids, namespace, initial_batch=100)

    new_docs = []
    new_ids = []
    for d, i in zip(docs, ids):
        if i not in existing_ids:
            new_docs.append(d)
            new_ids.append(i)

    return new_docs, new_ids


# ---------- Enhanced file loaders for multiple formats ----------

def _detect_encoding(file_path: Path) -> str:
    """Detect file encoding using chardet."""
    with open(file_path, 'rb') as f:
        raw_data = f.read(10000)  # Read first 10KB for detection
    result = chardet.detect(raw_data)
    return result.get('encoding', 'utf-8') or 'utf-8'


def _load_csv(path: Path) -> List[Document]:
    """Load CSV files with enhanced metadata and error handling."""
    encoding = _detect_encoding(path)
    try:
        loader = CSVLoader(
            file_path=str(path),
            encoding=encoding,
            csv_args={"delimiter": ",", "quoting": 0},
        )
        docs = loader.load()
        for i, d in enumerate(docs):
            d.metadata = d.metadata or {}
            d.metadata.update({
                "source": str(path),
                "filetype": "csv",
                "row_number": i + 1,
                "encoding": encoding
            })
        return docs
    except Exception as e:
        # Fallback with different encodings
        for enc in ("utf-8-sig", "cp1252", "latin-1"):
            try:
                loader = CSVLoader(
                    file_path=str(path),
                    encoding=enc,
                    csv_args={"delimiter": ",", "quoting": 0},
                )
                docs = loader.load()
                for i, d in enumerate(docs):
                    d.metadata = d.metadata or {}
                    d.metadata.update({
                        "source": str(path),
                        "filetype": "csv",
                        "row_number": i + 1,
                        "encoding": enc
                    })
                return docs
            except Exception:
                continue
        raise RuntimeError(f"Could not decode CSV file {path}") from e


def _load_pdf(path: Path) -> List[Document]:
    """Load PDF files with enhanced metadata."""
    loader = PyPDFLoader(str(path))
    docs = loader.load()
    for d in docs:
        d.metadata = d.metadata or {}
        d.metadata.update({
            "source": str(path),
            "filetype": "pdf",
            "total_pages": len(docs)
        })
    return docs


def _load_word_doc(path: Path) -> List[Document]:
    """Load Word documents (.docx, .doc)."""
    try:
        doc = DocxDocument(str(path))
        text_content = []
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text.strip())
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        full_text = "\n\n".join(text_content)
        
        if not full_text.strip():
            return []
            
        return [Document(
            page_content=full_text,
            metadata={
                "source": str(path),
                "filetype": "docx" if path.suffix.lower() == ".docx" else "doc",
                "paragraphs": len([p for p in doc.paragraphs if p.text.strip()]),
                "tables": len(doc.tables)
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process Word document {path}: {e}")


def _load_powerpoint(path: Path) -> List[Document]:
    """Load PowerPoint presentations (.pptx, .ppt)."""
    try:
        prs = Presentation(str(path))
        docs = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                content = "\n".join(slide_text)
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": str(path),
                        "filetype": "pptx" if path.suffix.lower() == ".pptx" else "ppt",
                        "slide": i + 1,
                        "total_slides": len(prs.slides)
                    }
                ))
        
        return docs
    except Exception as e:
        raise RuntimeError(f"Could not process PowerPoint file {path}: {e}")


def _load_html(path: Path) -> List[Document]:
    """Load HTML files."""
    encoding = _detect_encoding(path)
    try:
        with open(path, 'r', encoding=encoding) as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Extract text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        text = '\n'.join(line for line in lines if line)
        
        if not text.strip():
            return []
        
        # Extract title if available
        title = soup.find('title')
        title_text = title.get_text().strip() if title else ""
        
        return [Document(
            page_content=text,
            metadata={
                "source": str(path),
                "filetype": "html",
                "title": title_text,
                "encoding": encoding
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process HTML file {path}: {e}")


def _load_xml(path: Path) -> List[Document]:
    """Load XML files."""
    encoding = _detect_encoding(path)
    try:
        tree = ET.parse(str(path))
        root = tree.getroot()
        
        # Extract all text content
        text_content = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text_content.append(f"{elem.tag}: {elem.text.strip()}")
        
        if not text_content:
            return []
        
        full_text = "\n".join(text_content)
        
        return [Document(
            page_content=full_text,
            metadata={
                "source": str(path),
                "filetype": "xml",
                "root_tag": root.tag,
                "encoding": encoding
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process XML file {path}: {e}")


def _load_json(path: Path) -> List[Document]:
    """Load JSON files."""
    encoding = _detect_encoding(path)
    try:
        with open(path, 'r', encoding=encoding) as f:
            data = json.load(f)
        
        # Convert JSON to readable text
        if isinstance(data, dict):
            text_parts = []
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    text_parts.append(f"{key}: {json.dumps(value, indent=2)}")
                else:
                    text_parts.append(f"{key}: {value}")
            text = "\n".join(text_parts)
        elif isinstance(data, list):
            text = json.dumps(data, indent=2)
        else:
            text = str(data)
        
        return [Document(
            page_content=text,
            metadata={
                "source": str(path),
                "filetype": "json",
                "data_type": type(data).__name__,
                "encoding": encoding
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process JSON file {path}: {e}")


def _load_markdown(path: Path) -> List[Document]:
    """Load Markdown files."""
    encoding = _detect_encoding(path)
    try:
        with open(path, 'r', encoding=encoding) as f:
            content = f.read()
        
        # Convert markdown to HTML then extract text for better structure preservation
        html = markdown.markdown(content)
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text()
        
        # Clean up whitespace while preserving structure
        lines = (line.strip() for line in text.splitlines())
        text = '\n'.join(line for line in lines if line)
        
        return [Document(
            page_content=text,
            metadata={
                "source": str(path),
                "filetype": "markdown",
                "original_content": content[:500] + "..." if len(content) > 500 else content,
                "encoding": encoding
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process Markdown file {path}: {e}")


def _load_rtf(path: Path) -> List[Document]:
    """Load RTF files."""
    encoding = _detect_encoding(path)
    try:
        with open(path, 'r', encoding=encoding) as f:
            rtf_content = f.read()
        
        text = rtf_to_text(rtf_content)
        
        if not text.strip():
            return []
        
        return [Document(
            page_content=text,
            metadata={
                "source": str(path),
                "filetype": "rtf",
                "encoding": encoding
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process RTF file {path}: {e}")


def _load_code_file(path: Path) -> List[Document]:
    """Load code files with language-specific handling."""
    encoding = _detect_encoding(path)
    try:
        with open(path, 'r', encoding=encoding) as f:
            content = f.read()
        
        if not content.strip():
            return []
        
        language = CODE_EXTENSIONS.get(path.suffix.lower())
        
        return [Document(
            page_content=content,
            metadata={
                "source": str(path),
                "filetype": "code",
                "language": language.value if language else path.suffix.lower().lstrip('.'),
                "extension": path.suffix.lower(),
                "encoding": encoding,
                "lines": len(content.splitlines())
            }
        )]
    except Exception as e:
        raise RuntimeError(f"Could not process code file {path}: {e}")


def _load_text_like(path: Path) -> List[Document]:
    """Load plain text files with encoding detection."""
    encoding = _detect_encoding(path)
    try:
        loader = TextLoader(str(path), encoding=encoding)
        docs = loader.load()
        for d in docs:
            d.metadata = d.metadata or {}
            d.metadata.update({
                "source": str(path),
                "filetype": path.suffix.lstrip(".") or "txt",
                "encoding": encoding
            })
        return docs
    except Exception as e:
        # Fallback encodings
        for enc in ("utf-8-sig", "cp1252", "latin-1"):
            try:
                loader = TextLoader(str(path), encoding=enc)
                docs = loader.load()
                for d in docs:
                    d.metadata = d.metadata or {}
                    d.metadata.update({
                        "source": str(path),
                        "filetype": path.suffix.lstrip(".") or "txt",
                        "encoding": enc
                    })
                return docs
            except Exception:
                continue
        raise RuntimeError(f"Could not decode text file {path}") from e


def _docs_from_dataframe(df: pd.DataFrame, source: str, sheet: Optional[str]) -> List[Document]:
    """Convert DataFrame to Documents with enhanced metadata."""
    docs: List[Document] = []
    non_empty = df.fillna("")
    cols = list(non_empty.columns)
    
    for idx, row in non_empty.iterrows():
        parts = [f"{c}: {row[c]}" for c in cols if str(row[c]).strip() != ""]
        text = " | ".join(parts) if parts else ""
        if not text.strip():
            continue
        docs.append(
            Document(
                page_content=text,
                metadata={
                    "source": source,
                    "filetype": "excel",
                    "sheet": sheet,
                    "row_id": int(idx) if isinstance(idx, (int, float)) else str(idx),
                    "columns": cols,
                    "total_rows": len(df)
                },
            )
        )
    return docs


def _load_excel(path: Path) -> List[Document]:
    """Load Excel files with enhanced sheet handling."""
    docs: List[Document] = []
    try:
        xls = pd.read_excel(path, sheet_name=None)
        for sheet_name, df in xls.items():
            docs.extend(_docs_from_dataframe(df, source=str(path), sheet=sheet_name))
        return docs
    except Exception as e:
        raise RuntimeError(f"Could not process Excel file {path}: {e}")


def _load_any(folder: Path) -> List[Document]:
    """Load all supported file types from a folder."""
    docs: List[Document] = []
    
    # Define file type mappings
    loaders = {
        # Existing formats
        '.csv': _load_csv,
        '.pdf': _load_pdf,
        '.xlsx': _load_excel,
        '.xls': _load_excel,
        '.txt': _load_text_like,
        '.log': _load_text_like,
        
        # New formats
        '.docx': _load_word_doc,
        '.doc': _load_word_doc,
        '.pptx': _load_powerpoint,
        '.ppt': _load_powerpoint,
        '.html': _load_html,
        '.htm': _load_html,
        '.xml': _load_xml,
        '.json': _load_json,
        '.md': _load_markdown,
        '.markdown': _load_markdown,
        '.rtf': _load_rtf,
        
        # Code files
        '.py': _load_code_file,
        '.js': _load_code_file,
        '.jsx': _load_code_file,
        '.ts': _load_code_file,
        '.tsx': _load_code_file,
        '.java': _load_code_file,
        '.cpp': _load_code_file,
        '.c': _load_code_file,
        '.cs': _load_code_file,
        '.go': _load_code_file,
        '.php': _load_code_file,
        '.rb': _load_code_file,
        '.rs': _load_code_file,
        '.scala': _load_code_file,
        '.swift': _load_code_file,
        '.kt': _load_code_file,
        '.r': _load_code_file,
        '.sql': _load_code_file,
        '.css': _load_code_file,
        '.scss': _load_code_file,
        '.less': _load_code_file,
    }
    
    for p in sorted(folder.glob("*")):
        if not p.is_file():
            continue
            
        suffix = p.suffix.lower()
        loader = loaders.get(suffix)
        
        if loader:
            try:
                file_docs = loader(p)
                docs.extend(file_docs)
                print(f"Loaded {len(file_docs)} documents from {p.name}")
            except Exception as e:
                print(f"Warning: Could not load {p.name}: {e}")
        else:
            print(f"Skipping unsupported file type: {p.name}")
    
    return docs


# ---------- Enhanced text splitting strategies ----------

def _get_optimal_splitter(doc: Document) -> RecursiveCharacterTextSplitter:
    """Get optimal text splitter based on document type."""
    filetype = doc.metadata.get("filetype", "").lower()
    
    if filetype == "markdown":
        return MarkdownTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
        )
    elif filetype == "code":
        language = doc.metadata.get("language")
        # Check if language is supported by langchain Language enum
        supported_languages = [lang.value for lang in Language if hasattr(Language, lang.name)]
        if language and language in supported_languages:
            try:
                return PythonCodeTextSplitter.from_language(
                    language=Language(language),
                    chunk_size=800,
                    chunk_overlap=100,
                )
            except Exception:
                pass
    elif filetype in ["html", "xml"]:
        return RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=150,
            separators=["</div>", "</p>", "</section>", "\n\n", "\n", ".", " ", ""],
        )
    elif filetype in ["csv", "excel"]:
        # For structured data, use smaller chunks to keep rows together
        return RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=50,
            separators=["\n", " | ", ".", " ", ""],
        )
    elif filetype == "json":
        return RecursiveCharacterTextSplitter(
            chunk_size=1200,
            chunk_overlap=100,
            separators=["\n", "},", "]", ".", " ", ""],
        )
    
    # Default splitter for other formats
    return RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", "?", "!", " ", ""],
    )


def _split_docs(docs: List[Document]) -> List[Document]:
    """Split documents using optimal strategies for each file type."""
    all_splits = []
    
    # Group documents by type for efficient processing
    doc_groups = {}
    for doc in docs:
        filetype = doc.metadata.get("filetype", "default")
        if filetype not in doc_groups:
            doc_groups[filetype] = []
        doc_groups[filetype].append(doc)
    
    # Process each group with appropriate splitter
    for filetype, group_docs in doc_groups.items():
        if not group_docs:
            continue
            
        # Get splitter for this file type
        splitter = _get_optimal_splitter(group_docs[0])
        
        # Split documents
        splits = splitter.split_documents(group_docs)
        
        # Add chunk metadata
        for i, split in enumerate(splits):
            split.metadata = split.metadata or {}
            split.metadata.update({
                "chunk_index": i,
                "total_chunks": len(splits),
                "chunk_size": len(split.page_content)
            })
        
        all_splits.extend(splits)
    
    return all_splits


# ---------- Main ingestion function ----------

def ingest_folder(
    domain: Domain,
    data_dir: Path,
    namespace: Optional[str] = None,
) -> int:
    """Ingest all supported files in data_dir into Pinecone with enhanced processing."""
    print(f"Starting ingestion for domain: {domain}, directory: {data_dir}")
    
    # 1) Load all supported files
    docs = _load_any(data_dir)
    if not docs:
        print("No documents found to ingest")
        return 0
    
    print(f"Loaded {len(docs)} documents from {len(set(d.metadata.get('source') for d in docs))} files")
    
    # 2) Enrich metadata
    for d in docs:
        d.metadata = d.metadata or {}
        d.metadata.update({
            "domain": domain,
            "ingestion_timestamp": pd.Timestamp.now().isoformat(),
            "namespace": namespace or "default"
        })
    
    # 3) Split into chunks with optimal strategies
    splits = _split_docs(docs)
    print(f"Created {len(splits)} chunks after splitting")
    
    # 4) Filter out already-indexed docs
    index_name = _index_name(domain)
    new_splits, new_ids = _filter_new_docs(index_name, namespace, splits)
    if not new_splits:
        print("No new documents to index (all already exist)")
        return 0
    
    print(f"Found {len(new_splits)} new chunks to index")
    
    # 5) Embed and store
    embeddings = OpenAIEmbeddings(
        model=settings.OPENAI_EMBEDDING_MODEL,
        api_key=settings.OPENAI_API_KEY,
    )
    
    vectorstore = PineconeVectorStore(
        index_name=index_name,
        embedding=embeddings,
        namespace=namespace,
        pinecone_api_key=settings.PINECONE_API_KEY,
    )
    
    # Process in batches to avoid memory issues
    batch_size = 100
    total_indexed = 0
    
    for i in range(0, len(new_splits), batch_size):
        batch_docs = new_splits[i:i + batch_size]
        batch_ids = new_ids[i:i + batch_size]
        
        vectorstore.add_documents(batch_docs, ids=batch_ids)
        total_indexed += len(batch_docs)
        print(f"Indexed batch {i//batch_size + 1}: {len(batch_docs)} documents")
    
    print(f"Successfully indexed {total_indexed} new chunks")
    return total_indexed


# ---------- File upload support ----------

async def ingest_uploaded_file(
    domain: Domain,
    file_content: bytes,
    filename: str,
    namespace: Optional[str] = None,
    metadata: Optional[dict] = None,
) -> int:
    """Ingest a single uploaded file."""
    import tempfile
    import os
    from app.log import app_logger
    
    app_logger.info(f"Starting ingestion for file '{filename}' (size: {len(file_content)} bytes) in domain '{domain}'")
    
    # Create temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix) as tmp_file:
        tmp_file.write(file_content)
        tmp_path = Path(tmp_file.name)
    
    try:
        # Load the single file
        suffix = tmp_path.suffix.lower()
        
        # Map to appropriate loader
        loaders = {
            '.csv': _load_csv,
            '.pdf': _load_pdf,
            '.xlsx': _load_excel,
            '.xls': _load_excel,
            '.docx': _load_word_doc,
            '.doc': _load_word_doc,
            '.pptx': _load_powerpoint,
            '.ppt': _load_powerpoint,
            '.html': _load_html,
            '.htm': _load_html,
            '.xml': _load_xml,
            '.json': _load_json,
            '.md': _load_markdown,
            '.markdown': _load_markdown,
            '.rtf': _load_rtf,
            '.txt': _load_text_like,
        }
        
        # Add code file loaders
        for ext in CODE_EXTENSIONS:
            loaders[ext] = _load_code_file
        
        loader = loaders.get(suffix)
        if not loader:
            app_logger.error(f"Unsupported file type: {suffix} for file '{filename}'")
            raise ValueError(f"Unsupported file type: {suffix}")
        
        app_logger.info(f"Loading file '{filename}' using {loader.__name__}")
        docs = loader(tmp_path)
        app_logger.info(f"Loaded {len(docs)} documents from file '{filename}'")
        if not docs:
            from app.log import app_logger
            app_logger.warning(f"File '{filename}' loaded but produced no documents (file may be empty or unreadable)")
            return 0
        
        # Update metadata with original filename and additional metadata
        for d in docs:
            d.metadata = d.metadata or {}
            d.metadata.update({
                "source": filename,
                "domain": domain,
                "ingestion_timestamp": pd.Timestamp.now().isoformat(),
                "namespace": namespace or "default",
                "filename": filename,
                "file_type": Path(filename).suffix.lower(),
                "document_id": f"doc_{hashlib.md5(filename.encode()).hexdigest()[:8]}"
            })
            
            # Add any additional metadata passed in
            if metadata:
                d.metadata.update(metadata)
        
        # Split documents
        app_logger.info(f"Splitting {len(docs)} documents into chunks for file '{filename}'")
        splits = _split_docs(docs)
        app_logger.info(f"Created {len(splits)} chunks from file '{filename}'")
        
        # Filter and ingest
        index_name = _index_name(domain)
        app_logger.info(f"Checking for existing chunks in index '{index_name}' (namespace: {namespace})")
        new_splits, new_ids = _filter_new_docs(index_name, namespace, splits)
        
        if not new_splits:
            app_logger.warning(
                f"No new chunks to index for file '{filename}' in domain '{domain}' "
                f"(namespace: {namespace}). All {len(splits)} chunks already exist. "
                f"This may mean the file was already uploaded."
            )
            return 0
        
        app_logger.info(f"Found {len(new_splits)} new chunks to index (out of {len(splits)} total) for file '{filename}'")
        
        # Embed and store
        app_logger.info(f"Creating embeddings and storing {len(new_splits)} chunks for file '{filename}'")
        embeddings = OpenAIEmbeddings(
            model=settings.OPENAI_EMBEDDING_MODEL,
            api_key=settings.OPENAI_API_KEY,
        )
        
        vectorstore = PineconeVectorStore(
            index_name=index_name,
            embedding=embeddings,
            namespace=namespace,
            pinecone_api_key=settings.PINECONE_API_KEY,
        )
        
        vectorstore.add_documents(new_splits, ids=new_ids)
        app_logger.info(f"Successfully indexed {len(new_splits)} chunks for file '{filename}'")
        return len(new_splits)
        
    except Exception as e:
        from app.log import app_logger
        app_logger.error(f"Error ingesting file '{filename}': {str(e)}", exc_info=True)
        raise
    finally:
        # Clean up temporary file
        if tmp_path.exists():
            try:
                os.unlink(tmp_path)
            except Exception as e:
                from app.log import app_logger
                app_logger.warning(f"Could not delete temporary file {tmp_path}: {e}")